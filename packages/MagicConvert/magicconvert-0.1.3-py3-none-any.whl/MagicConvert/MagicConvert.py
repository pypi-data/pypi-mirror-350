import copy
import html
import mimetypes
import os
import platform
import re
import tempfile
import traceback
from typing import Any, Dict, List, Optional, Union
from urllib.parse import quote, unquote, urlparse, urlunparse
from warnings import warn, resetwarnings, catch_warnings
import mammoth
import markdownify
import pandas as pd
import pdfminer
import pdfminer.high_level
import pptx
import puremagic
import requests
from bs4 import BeautifulSoup
from charset_normalizer import from_path
from PIL import Image
import pytesseract
import warnings
warnings.filterwarnings("ignore")

def setup_tesseract():
    """
    Automatically configure Tesseract for Windows and Linux systems
    """
    system = platform.system().lower()
    
    if system == "windows":
        # Common Tesseract installation paths on Windows
        possible_paths = [
            r"C:\Program Files\Tesseract-OCR\tesseract.exe",
            r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
            r"C:\Users\%USERNAME%\AppData\Local\Programs\Tesseract-OCR\tesseract.exe",
            r"C:\tesseract\tesseract.exe",
        ]
        
        # Expand environment variables in paths
        expanded_paths = []
        for path in possible_paths:
            expanded_path = os.path.expandvars(path)
            expanded_paths.append(expanded_path)
        
        # Find the first existing Tesseract executable
        for path in expanded_paths:
            if os.path.isfile(path):
                pytesseract.pytesseract.tesseract_cmd = path
                print(f"Tesseract found at: {path}")
                return path
        
        # If not found, try to find it in PATH
        import shutil
        tesseract_path = shutil.which("tesseract")
        if tesseract_path:
            pytesseract.pytesseract.tesseract_cmd = tesseract_path
            print(f"Tesseract found in PATH: {tesseract_path}")
            return tesseract_path
        
        # If still not found, provide helpful error message
        raise FileNotFoundError(
            "Tesseract not found on Windows. Please install Tesseract OCR from: "
            "https://github.com/UB-Mannheim/tesseract/wiki\n"
            "Or ensure it's installed in one of these locations:\n" + 
            "\n".join(f"  - {path}" for path in expanded_paths)
        )
    
    elif system == "linux":
        # Common paths for Linux
        possible_paths = [
            "/usr/bin/tesseract",
            "/usr/local/bin/tesseract",
            "/opt/tesseract/bin/tesseract",
        ]
        
        for path in possible_paths:
            if os.path.isfile(path):
                pytesseract.pytesseract.tesseract_cmd = path
                print(f"Tesseract found at: {path}")
                return path
        
        # Try to find in PATH
        import shutil
        tesseract_path = shutil.which("tesseract")
        if tesseract_path:
            pytesseract.pytesseract.tesseract_cmd = tesseract_path
            print(f"Tesseract found in PATH: {tesseract_path}")
            return tesseract_path
        
        raise FileNotFoundError(
            "Tesseract not found on Linux. Please install it using:\n"
            "Ubuntu/Debian: sudo apt-get install tesseract-ocr\n"
            "CentOS/RHEL: sudo yum install tesseract\n"
            "Fedora: sudo dnf install tesseract"
        )
    
    elif system == "darwin":  # macOS
        possible_paths = [
            "/usr/local/bin/tesseract",
            "/opt/homebrew/bin/tesseract",
        ]
        
        for path in possible_paths:
            if os.path.isfile(path):
                pytesseract.pytesseract.tesseract_cmd = path
                print(f"Tesseract found at: {path}")
                return path
        
        import shutil
        tesseract_path = shutil.which("tesseract")
        if tesseract_path:
            pytesseract.pytesseract.tesseract_cmd = tesseract_path
            print(f"Tesseract found in PATH: {tesseract_path}")
            return tesseract_path
        
        raise FileNotFoundError(
            "Tesseract not found on macOS. Please install it using:\n"
            "brew install tesseract"
        )
    
    else:
        # For other systems, try to find in PATH
        import shutil
        tesseract_path = shutil.which("tesseract")
        if tesseract_path:
            pytesseract.pytesseract.tesseract_cmd = tesseract_path
            print(f"Tesseract found in PATH: {tesseract_path}")
            return tesseract_path
        
        raise FileNotFoundError(f"Tesseract not found on {system}. Please install Tesseract OCR.")

# Setup Tesseract automatically
try:
    setup_tesseract()
except FileNotFoundError as e:
    print(f"Warning: {e}")
    print("OCR functionality will not be available until Tesseract is properly installed.")

class OCRReader:
    def __init__(self, tesseract_cmd: Optional[str] = None, config: Optional[Dict] = None):
        # Use provided tesseract_cmd or keep the auto-detected one
        if tesseract_cmd:
            pytesseract.pytesseract.tesseract_cmd = tesseract_cmd
        elif not hasattr(pytesseract.pytesseract, 'tesseract_cmd') or not pytesseract.pytesseract.tesseract_cmd:
            # If no tesseract_cmd is set, try to set it up
            try:
                setup_tesseract()
            except FileNotFoundError:
                pass  # Will be handled when trying to use OCR
        
        self.config = config or {}

    def read_text_from_image(self, image: Image.Image) -> str:
        try:
            # Verify Tesseract is available before attempting OCR
            if not pytesseract.pytesseract.tesseract_cmd:
                raise Exception("Tesseract is not configured. Please install Tesseract OCR.")
            
            # Test if Tesseract is accessible
            if not os.path.isfile(pytesseract.pytesseract.tesseract_cmd):
                raise Exception(f"Tesseract not found at: {pytesseract.pytesseract.tesseract_cmd}")
            
            text = pytesseract.image_to_string(image, **self.config)
            return text.strip()
        except Exception as e:
            raise Exception(f"Error processing image: {str(e)}")


class _CustomMarkdownify(markdownify.MarkdownConverter):
    def __init__(self, **options: Any):
        options["heading_style"] = options.get("heading_style", markdownify.ATX)
        super().__init__(**options)

    def convert_a(self, el: Any, text: str, *args, **kwargs):
        prefix, suffix, text = markdownify.chomp(text)
        if not text:
            return ""
        href = el.get("href")
        title = el.get("title")
        if href:
            try:
                parsed_url = urlparse(href)  # type: ignore
                if parsed_url.scheme and parsed_url.scheme.lower() not in ["http", "https", "file"]:  # type: ignore
                    return "%s%s%s" % (prefix, text, suffix)
                href = urlunparse(parsed_url._replace(path=quote(unquote(parsed_url.path))))  # type: ignore
            except ValueError:
                return "%s%s%s" % (prefix, text, suffix)
        if (
            self.options["autolinks"]
            and text.replace(r"\_", "_") == href
            and not title
            and not self.options["default_title"]
        ):
            return "<%s>" % href
        if self.options["default_title"] and not title:
            title = href
        title_part = ' "%s"' % title.replace('"', r"\"") if title else ""
        return (
            "%s[%s](%s%s)%s" % (prefix, text, href, title_part, suffix)
            if href
            else text
        )

    def convert_hn(self, n: int, el: Any, text: str, convert_as_inline: bool) -> str:
        if not convert_as_inline:
            if not re.search(r"^\n", text):
                return "\n" + super().convert_hn(n, el, text, convert_as_inline)  # type: ignore

        return super().convert_hn(n, el, text, convert_as_inline)  # type: ignore

    def convert_img(self, el: Any, text: str, *args, **kwargs) -> str:
        # Handle both old and new calling patterns
        convert_as_inline = kwargs.get('convert_as_inline', False)
        if len(args) > 0:
            convert_as_inline = args[0]
            
        alt = el.attrs.get("alt", None) or ""
        src = el.attrs.get("src", None) or ""
        title = el.attrs.get("title", None) or ""
        title_part = ' "%s"' % title.replace('"', r"\"") if title else ""
        if (
            convert_as_inline
            and el.parent.name not in self.options["keep_inline_images_in"]
        ):
            return alt
        if src.startswith("data:"):
            src = src.split(",")[0] + "..."

        return "![%s](%s%s)" % (alt, src, title_part)

    def convert_soup(self, soup: Any) -> str:
        return super().convert_soup(soup)


class DocumentConverterResult:
    def __init__(self, title: Union[str, None] = None, text_content: str = ""):
        self.title: Union[str, None] = title
        self.text_content: str = text_content


class DocumentConverter:
    def convert(self, local_path: str, **kwargs: Any) -> Union[None, DocumentConverterResult]:
        raise NotImplementedError()

    def supports_extension(self, ext: str) -> bool:
        """Return True if this converter supports the given extension."""
        raise NotImplementedError()


class PlainTextConverter(DocumentConverter):
    def convert(
        self, local_path: str, **kwargs: Any
    ) -> Union[None, DocumentConverterResult]:
        content_type, _ = mimetypes.guess_type(
            "__placeholder" + kwargs.get("file_extension", "")
        )
        if content_type is None:
            return None
        elif "text/" not in content_type.lower():
            return None

        text_content = str(from_path(local_path).best())
        return DocumentConverterResult(
            title=None,
            text_content=text_content,
        )


class HtmlConverter(DocumentConverter):
    def convert(
        self, local_path: str, **kwargs: Any
    ) -> Union[None, DocumentConverterResult]:
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".html", ".htm"]:
            return None

        result = None
        with open(local_path, "rt", encoding="utf-8") as fh:
            result = self._convert(fh.read())

        return result

    def _convert(self, html_content: str) -> Union[None, DocumentConverterResult]:
        soup = BeautifulSoup(html_content, "html.parser")
        for script in soup(["script", "style"]):
            script.extract()
        body_elm = soup.find("body")
        webpage_text = ""
        if body_elm:
            webpage_text = _CustomMarkdownify().convert_soup(body_elm)
        else:
            webpage_text = _CustomMarkdownify().convert_soup(soup)

        assert isinstance(webpage_text, str)

        return DocumentConverterResult(
            title=None if soup.title is None else soup.title.string,
            text_content=webpage_text,
        )


class PdfConverter(DocumentConverter):
    def supports_extension(self, ext: str) -> bool:
        return ext.lower() == '.pdf'

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".pdf":
            return None
        return DocumentConverterResult(
            title=None,
            text_content=pdfminer.high_level.extract_text(local_path),
        )


class DocxConverter(HtmlConverter):
    def supports_extension(self, ext: str) -> bool:
        return ext.lower() == '.docx'

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".docx":
            return None
        result = None
        with open(local_path, "rb") as docx_file:
            style_map = kwargs.get("style_map", None)
            result = mammoth.convert_to_html(docx_file, style_map=style_map)
            html_content = result.value
            result = self._convert(html_content)
        return result


class XlsxConverter(HtmlConverter):
    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".xlsx":
            return None

        sheets = pd.read_excel(local_path, sheet_name=None)
        md_content = ""
        for s in sheets:
            md_content += f"## {s}\n"
            html_content = sheets[s].to_html(index=False)
            md_content += self._convert(html_content).text_content.strip() + "\n\n"

        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )


class PptxConverter(HtmlConverter):
    def supports_extension(self, ext: str) -> bool:
        return ext.lower() == '.pptx'

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".pptx":
            return None
        md_content = ""
        presentation = pptx.Presentation(local_path)
        slide_num = 0
        for slide in presentation.slides:
            slide_num += 1

            md_content += f"\n\n<!-- Slide number: {slide_num} -->\n"

            title = slide.shapes.title
            for shape in slide.shapes:
                if self._is_picture(shape):
                    alt_text = ""
                    try:
                        alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
                    except Exception:
                        pass
                    filename = re.sub(r"\W", "", shape.name) + ".jpg"
                    md_content += (
                        "\n!["
                        + (alt_text if alt_text else shape.name)
                        + "]("
                        + filename
                        + ")\n"
                    )

                # Tables
                if self._is_table(shape):
                    html_table = "<html><body><table>"
                    first_row = True
                    for row in shape.table.rows:
                        html_table += "<tr>"
                        for cell in row.cells:
                            if first_row:
                                html_table += "<th>" + html.escape(cell.text) + "</th>"
                            else:
                                html_table += "<td>" + html.escape(cell.text) + "</td>"
                        html_table += "</tr>"
                        first_row = False
                    html_table += "</table></body></html>"
                    md_content += (
                        "\n" + self._convert(html_table).text_content.strip() + "\n"
                    )
                if shape.has_chart:
                    md_content += self._convert_chart_to_markdown(shape.chart)
                elif shape.has_text_frame:
                    if shape == title:
                        md_content += "# " + shape.text.lstrip() + "\n"
                    else:
                        md_content += shape.text + "\n"

            md_content = md_content.strip()

            if slide.has_notes_slide:
                md_content += "\n\n### Notes:\n"
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame is not None:
                    md_content += notes_frame.text
                md_content = md_content.strip()

        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )

    def _is_picture(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER:
            if hasattr(shape, "image"):
                return True
        return False

    def _is_table(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE:
            return True
        return False

    def _convert_chart_to_markdown(self, chart):
        md = "\n\n### Chart"
        if chart.has_title:
            md += f": {chart.chart_title.text_frame.text}"
        md += "\n\n"
        data = []
        category_names = [c.label for c in chart.plots[0].categories]
        series_names = [s.name for s in chart.series]
        data.append(["Category"] + series_names)

        for idx, category in enumerate(category_names):
            row = [category]
            for series in chart.series:
                row.append(series.values[idx])
            data.append(row)

        markdown_table = []
        for row in data:
            markdown_table.append("| " + " | ".join(map(str, row)) + " |")
        header = markdown_table[0]
        separator = "|" + "|".join(["---"] * len(data[0])) + "|"
        return md + "\n".join([header, separator] + markdown_table[1:])


class FileConversionException(BaseException):
    pass


class UnsupportedFormatException(BaseException):
    pass


class ImageConverter(DocumentConverter):
    def __init__(self, ocr_reader: Optional[OCRReader] = None):
        self.ocr_reader = ocr_reader or OCRReader()

    def convert(self, local_path: str, **kwargs: Any) -> Union[None, DocumentConverterResult]:
        extension = kwargs.get("file_extension", "").lower()
        if extension not in ['.png', '.jpg', '.jpeg', '.tiff', '.bmp']:
            return None

        try:
            image = Image.open(local_path)
            text_content = self.ocr_reader.read_text_from_image(image)
            markdown_content = self._convert_to_markdown_structure(text_content)
            return DocumentConverterResult(
                title=None,
                text_content=markdown_content
            )
        except Exception as e:
            raise FileConversionException(f"Failed to process image: {str(e)}")

    def _convert_to_markdown_structure(self, text_content: str) -> str:
        lines = text_content.split('\n')
        markdown = []
        current_table = []
        in_table = False

        i = 0
        while i < len(lines):
            line = lines[i].strip()
            next_line = lines[i + 1].strip() if i + 1 < len(lines) else ""

            if not line:
                if in_table:
                    markdown.extend(self._format_table(current_table))
                    current_table = []
                    in_table = False
                markdown.append("")
                i += 1
                continue

            header_level = self._detect_header_level(line, next_line)
            if header_level:
                if in_table:
                    markdown.extend(self._format_table(current_table))
                    current_table = []
                    in_table = False
                markdown.append(f"{'#' * header_level} {line}")
                i += 2 if header_level > 0 and next_line and set(next_line) in [set('='), set('-')] else 1
                continue

            list_format = self._detect_list_format(line)
            if list_format:
                if in_table:
                    markdown.extend(self._format_table(current_table))
                    current_table = []
                    in_table = False
                markdown.append(list_format)
                i += 1
                continue
            if self._is_likely_table_row(line):
                in_table = True
                current_table.append(line)
                i += 1
                continue
            if in_table:
                markdown.extend(self._format_table(current_table))
                current_table = []
                in_table = False
            line = self._format_emphasis(line)

            markdown.append(line)
            i += 1
        if current_table:
            markdown.extend(self._format_table(current_table))

        return "\n\n".join([l for l in markdown if l])

    def _detect_header_level(self, line: str, next_line: str) -> int:
        if line.startswith('#'):
            return len(re.match(r'^#+', line).group())
        if next_line:
            if set(next_line) == set('='):
                return 1
            if set(next_line) == set('-'):
                return 2
        if len(line) <= 100 and line.strip():
            words = line.split()
            if all(word[0].isupper() for word in words if word):
                return 1
            if line[0].isupper() and len(words) <= 10:
                return 2

        return 0

    def _detect_list_format(self, line: str) -> Optional[str]:
        bullet_points = ['-', '•', '*', '○', '►', '·']
        for bullet in bullet_points:
            if line.lstrip().startswith(bullet):
                content = line.lstrip()[1:].strip()
                return f"- {content}"

        if re.match(r'^\d+[\.\)]', line):
            content = re.sub(r'^\d+[\.\)]', '', line).strip()
            return f"1. {content}"

        return None

    def _is_likely_table_row(self, line: str) -> bool:
        parts = [p for p in re.split(r'\s{2,}', line) if p.strip()]
        if len(parts) >= 2:
            lengths = [len(p) for p in parts]
            avg_length = sum(lengths) / len(lengths)
            if all(abs(l - avg_length) <= 5 for l in lengths):
                return True
        return False

    def _format_table(self, table_rows: List[str]) -> List[str]:
        if not table_rows:
            return []
        split_rows = [re.split(r'\s{2,}', row.strip()) for row in table_rows]
        max_cols = max(len(row) for row in split_rows)
        normalized_rows = []
        for row in split_rows:
            while len(row) < max_cols:
                row.append('')
            normalized_rows.append(row)
        col_widths = []
        for col in range(max_cols):
            width = max(len(row[col]) for row in normalized_rows)
            col_widths.append(width)
        markdown_table = []

        header = "| " + " | ".join(cell.ljust(width) for cell, width in zip(normalized_rows[0], col_widths)) + " |"
        markdown_table.append(header)

        separator = "|" + "|".join("-" * (width + 2) for width in col_widths) + "|"
        markdown_table.append(separator)

        for row in normalized_rows[1:]:
            formatted_row = "| " + " | ".join(cell.ljust(width) for cell, width in zip(row, col_widths)) + " |"
            markdown_table.append(formatted_row)

        return markdown_table

    def _format_emphasis(self, text: str) -> str:
        text = re.sub(r'\b([A-Z]{2,})\b', r'**\1**', text)
        text = re.sub(r'[_/](.*?)[_/]', r'*\1*', text)
        return text

class MagicConvert:
    def __init__(
            self,
            requests_session: Optional[requests.Session] = None,
            style_map: Optional[str] = None,
    ):
        if requests_session is None:
            self._requests_session = requests.Session()
        else:
            self._requests_session = requests_session

        self._style_map = style_map
        self._page_converters: List[DocumentConverter] = []

        ocr_reader = OCRReader()

        self.register_page_converter(ImageConverter(ocr_reader))
        self.register_page_converter(PlainTextConverter())
        self.register_page_converter(HtmlConverter())
        self.register_page_converter(DocxConverter())
        self.register_page_converter(XlsxConverter())
        self.register_page_converter(PptxConverter())
        self.register_page_converter(PdfConverter())

    def magic(
        self, source: Union[str, requests.Response], **kwargs: Any
    ) -> DocumentConverterResult:
        if isinstance(source, str):
            if (
                source.startswith("http://")
                or source.startswith("https://")
                or source.startswith("file://")
            ):
                return self.convert_url(source, **kwargs)
            else:
                return self.convert_local(source, **kwargs)
        elif isinstance(source, requests.Response):
            return self.convert_response(source, **kwargs)

    def convert_local(
        self, path: str, **kwargs: Any
    ) -> DocumentConverterResult:
        ext = kwargs.get("file_extension")
        extensions = [ext] if ext is not None else []
        base, ext = os.path.splitext(path)
        self._append_ext(extensions, ext)

        for g in self._guess_ext_magic(path):
            self._append_ext(extensions, g)
        return self._convert(path, extensions, **kwargs)

    def convert_stream(
        self, stream: Any, **kwargs: Any
    ) -> DocumentConverterResult:
        ext = kwargs.get("file_extension")
        extensions = [ext] if ext is not None else []
        handle, temp_path = tempfile.mkstemp()
        fh = os.fdopen(handle, "wb")
        result = None
        try:
            content = stream.read()
            if isinstance(content, str):
                fh.write(content.encode("utf-8"))
            else:
                fh.write(content)
            fh.close()
            for g in self._guess_ext_magic(temp_path):
                self._append_ext(extensions, g)
            result = self._convert(temp_path, extensions, **kwargs)
        finally:
            try:
                fh.close()
            except Exception:
                pass
            os.unlink(temp_path)

        return result

    def convert_url(
        self, url: str, **kwargs: Any
    ) -> DocumentConverterResult:
        response = self._requests_session.get(url, stream=True)
        response.raise_for_status()
        return self.convert_response(response, **kwargs)

    def convert_response(
        self, response: requests.Response, **kwargs: Any
    ) -> DocumentConverterResult:
        ext = kwargs.get("file_extension")
        extensions = [ext] if ext is not None else []
        content_type = response.headers.get("content-type", "").split(";")[0]
        self._append_ext(extensions, mimetypes.guess_extension(content_type))
        content_disposition = response.headers.get("content-disposition", "")
        m = re.search(r"filename=([^;]+)", content_disposition)
        if m:
            base, ext = os.path.splitext(m.group(1).strip("\"'"))
            self._append_ext(extensions, ext)
        base, ext = os.path.splitext(urlparse(response.url).path)
        self._append_ext(extensions, ext)
        handle, temp_path = tempfile.mkstemp()
        fh = os.fdopen(handle, "wb")
        result = None
        try:
            for chunk in response.iter_content(chunk_size=512):
                fh.write(chunk)
            fh.close()
            for g in self._guess_ext_magic(temp_path):
                self._append_ext(extensions, g)

            result = self._convert(temp_path, extensions, url=response.url, **kwargs)
        finally:
            try:
                fh.close()
            except Exception:
                pass
            os.unlink(temp_path)

        return result

    def _convert(
    self, local_path: str, extensions: List[Union[str, None]], **kwargs
    ) -> DocumentConverterResult:
        error_trace = ""
        for ext in extensions + [None]:
            for converter in self._page_converters:
                _kwargs = copy.deepcopy(kwargs)
                if ext is None:
                    if "file_extension" in _kwargs:
                        del _kwargs["file_extension"]
                else:
                    _kwargs.update({"file_extension": ext})
    
                _kwargs["_parent_converters"] = self._page_converters
                if "style_map" not in _kwargs and self._style_map is not None:
                    _kwargs["style_map"] = self._style_map
    
                try:
                    res = converter.convert(local_path, **_kwargs)
                    if res is not None:
                        res.text_content = "\n".join(
                            [line.rstrip() for line in re.split(r"\r?\n", res.text_content)]
                        )
                        res.text_content = re.sub(r"\n{3,}", "\n\n", res.text_content)
                        return res
                except Exception as e:
                    # If this converter supports the extension and fails, raise the exception
                    if ext is not None and converter.supports_extension(ext):
                        raise FileConversionException(
                            f"Could not convert '{local_path}' to Markdown with {converter.__class__.__name__} "
                            f"for extension '{ext}'. The following error occurred:\n\n{traceback.format_exc()}"
                        )
                    # Otherwise, store the error and continue
                    error_trace = ("\n\n" + traceback.format_exc()).strip()
    
        if len(error_trace) > 0:
            raise FileConversionException(
                f"Could not convert '{local_path}' to Markdown. File type was recognized as {extensions}. "
                f"While converting the file, the following error was encountered:\n\n{error_trace}"
            )
        raise UnsupportedFormatException(
            f"Could not convert '{local_path}' to Markdown. The formats {extensions} are not supported."
        )

    def _append_ext(self, extensions, ext):
        if ext is None:
            return
        ext = ext.strip()
        if ext == "":
            return
        extensions.append(ext)

    def _guess_ext_magic(self, path):
        try:
            guesses = puremagic.magic_file(path)
            extensions = list()
            for g in guesses:
                ext = g.extension.strip()
                if len(ext) > 0:
                    if not ext.startswith("."):
                        ext = "." + ext
                    if ext not in extensions:
                        extensions.append(ext)
            return extensions
        except FileNotFoundError:
            pass
        except IsADirectoryError:
            pass
        except PermissionError:
            pass
        return []

    def register_page_converter(self, converter: DocumentConverter) -> None:
        self._page_converters.insert(0, converter)